"""PowerPoint comparison report builder using native Excel charts."""

import os
from math import ceil

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_AXIS_CROSSES, XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE, XL_TICK_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

from dashboard.constants import SELECTION_LABELS

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_BAR_TOP = Inches(0.25)
TITLE_BAR_H = Inches(0.75)
TITLE_BAR_COLOR = RGBColor(0x2a, 0x2a, 0x2a)  # Dark panel background

CHART_TOP = Inches(1.3)
CHART_H = Inches(5.3)
CHART_MARGIN = Inches(0.2)
MAX_CHARTS_PER_SLIDE = 3

SERIES_COLORS = [
    RGBColor(0x3b, 0x82, 0xf6),  # Blue
    RGBColor(0xef, 0x44, 0x44),  # Red
    RGBColor(0x22, 0xc5, 0x5e),  # Green
    RGBColor(0xf5, 0x9e, 0x0b),  # Amber
    RGBColor(0x8b, 0x5c, 0xf6),  # Purple
    RGBColor(0xec, 0x48, 0x99),  # Pink
    RGBColor(0x14, 0xb8, 0xa6),  # Teal
    RGBColor(0xf9, 0x73, 0x16),  # Orange
]

# Branding & visual polish – dark theme matching desktop app
ACCENT_COLOR = RGBColor(0x3b, 0x82, 0xf6)       # Blue accent (matches app accent)
SLIDE_BG = RGBColor(0x1f, 0x1f, 0x1f)           # Dark background (matches app bg_primary)
FOOTER_TEXT = RGBColor(0xa0, 0xa0, 0xa0)         # Secondary text (matches app text_secondary)
FOOTER_LINE = RGBColor(0x3d, 0x3d, 0x3d)        # Border separator (matches app border)
TABLE_ROW_ALT = RGBColor(0x2a, 0x2a, 0x2a)      # Alternating row (matches app bg_secondary)
TABLE_ROW_WHITE = RGBColor(0x35, 0x35, 0x35)     # Base row (matches app bg_tertiary)
TEXT_PRIMARY = RGBColor(0xe0, 0xe0, 0xe0)        # Primary text color
TEXT_SECONDARY = RGBColor(0xa0, 0xa0, 0xa0)      # Secondary text color
TEXT_MUTED = RGBColor(0x70, 0x70, 0x70)          # Muted text color

# Mapping from comparison category column to RawFile property
CATEGORY_PROP_MAP = {
    'Wafer': 'wafer_type',
    'Pad': 'pad_type',
    'Slurry': 'slurry_type',
    'Conditioner': 'conditioner_disk_type',
}

# Time-series metrics for appendix slides
# Each entry: (chart_title, [(column_name, series_label), ...])
TIME_TRACE_METRICS = [
    ('COF', [('COF', 'COF')]),
    ('Forces (lbf)', [('Fy Total (lbf)', 'Fy (Shear)'), ('Fz Total (lbf)', 'Fz (Down)')]),
    ('IR Temperature (\u00b0C)', [('IR Temperature', 'IR Temperature')]),
]

# Unified correlation graph definitions for PPTX reports
# (chart_id, x_col, y_col, x_label, y_label, transform_x, transform_y)
# First 7 are primary (default on), rest are secondary (default off)
PPTX_CORR_GRAPHS = [
    ('stribeck',     'Sommerfeld', 'COF',          'Pseudo-Sommerfeld Number', 'Mean COF',                       'log_axis', 'log_axis'),
    ('temp-cofpv',   'COF.P.V',   'Mean Temp',    'COF\u00b7P\u00b7V (Pa\u00b7m/s)',  'Mean Pad Temp (\u00b0C)',                 None, None),
    ('var-fy-pv',    'P.V',       'Var Fy',       'P\u00b7V (Pa\u00b7m/s)',            'Variance of Shear Force (lbf\u00b2)',     None, None),
    ('var-fz-pv',    'P.V',       'Var Fz',       'P\u00b7V (Pa\u00b7m/s)',            'Variance of Normal Force (lbf\u00b2)',    None, None),
    ('rr-cofpv',     'COF.P.V',   'Removal Rate', 'COF\u00b7P\u00b7V (Pa\u00b7m/s)',  'Removal Rate (\u00c5/min)',               None, None),
    ('directivity',  'P.V',       '_directivity', 'P\u00b7V (Pa\u00b7m/s)',            'Directivity',                            None, 'log_axis'),
    ('wiwrrnu-pv',   'P.V',       'WIWNU',        'P\u00b7V (Pa\u00b7m/s)',            'WIWRRNU (%)',                             None, None),
    ('cof-pv',       'P.V',       'COF',          'P\u00b7V (Pa\u00b7m/s)',            'Mean COF',                                None, None),
    ('temp-pv',      'P.V',       'Mean Temp',    'P\u00b7V (Pa\u00b7m/s)',            'Mean Pad Temp (\u00b0C)',                 None, None),
    ('arrhenius',    'Mean Temp', 'Removal Rate', '1/T (1/K)',                         'ln(RR)',                                  'inv_kelvin', 'ln'),
    ('preston',      'P.V',       'Removal Rate', 'P\u00b7V (Pa\u00b7m/s)',            'Mean Removal Rate (\u00c5/min)',          None, None),
]


class PptxBuilder:
    """Builds a PowerPoint comparison report with native editable charts."""

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    @staticmethod
    def _count_secondary_values(report, secondary):
        """Count unique non-empty values of the secondary category."""
        prop = CATEGORY_PROP_MAP.get(secondary)
        if not prop:
            return 1
        seen = set()
        for f in report.files:
            val = getattr(f, prop, None)
            if val and str(val).strip():
                seen.add(str(val).strip())
        return max(len(seen), 1)

    @staticmethod
    def count_slides(report, config):
        """Pre-compute exact slide count without building anything."""
        num_groups = PptxBuilder._count_groups(report, config['comparison_category'])
        n_metrics = len(config.get('selected_metrics', []))
        n_corr = len(config.get('selected_correlations', []))
        n_sec = PptxBuilder._count_secondary_values(
            report, config.get('secondary_category', 'Wafer'))
        groups_per_slide = max(1, ceil(num_groups / MAX_CHARTS_PER_SLIDE))

        slides = 2  # title + TOC
        # Executive summary: one slide per batch of 3 bar charts
        slides += max(1, ceil(n_metrics / MAX_CHARTS_PER_SLIDE))
        # Correlations: one slide per (graph × secondary value)
        slides += n_corr * n_sec * groups_per_slide
        # Summary tables: one per (group × secondary value)
        if config.get('include_summary_tables'):
            slides += num_groups * n_sec
        # Time traces
        if config.get('include_time_traces'):
            slides += len(report.files)
        return slides

    @classmethod
    def build(cls, report, config, project_dir, progress_cb):
        """Build and save the PowerPoint file. Returns the file path.

        Args:
            progress_cb: callable(label: str) called after each slide is created.
        """
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        category = config['comparison_category']
        prop = CATEGORY_PROP_MAP[category]
        project_name = os.path.basename(project_dir)

        # Resolve logo path (project root)
        app_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        logo_path = os.path.join(app_root, 'logo_inverted.png')
        if not os.path.isfile(logo_path):
            logo_path = None

        # Build data
        df = cls._build_summary_df(report, category, prop)
        groups = cls._get_groups(df, category)
        group_names = sorted(groups.keys())

        # Determine PSI series within groups
        psi_values = cls._get_psi_values(df)

        # ---- Slides ----
        cls._add_title_slide(prs, project_name, category, config, logo_path)
        progress_cb("Creating title slide...")

        cls._add_toc_slide(prs, config, group_names)
        progress_cb("Creating table of contents...")

        # Executive summary bar charts
        metrics = config.get('selected_metrics', [])
        secondary = config.get('secondary_category', 'Wafer')
        cls._add_executive_summary(prs, df, groups, group_names, metrics, category, secondary, progress_cb)

        # Correlation graphs (split by secondary)
        corr_configs = config.get('selected_correlations', [])
        cls._add_correlation_slides(prs, df, groups, group_names, corr_configs, psi_values, category, secondary, progress_cb)

        # Summary tables
        if config.get('include_summary_tables'):
            summary_features = config.get('summary_table_features', [])
            cls._add_summary_tables(prs, df, groups, group_names, summary_features, category, secondary, progress_cb)

        # Time traces appendix
        if config.get('include_time_traces'):
            cof_y_range = (config.get('tt_cof_y_min'), config.get('tt_cof_y_max'))
            cls._add_time_traces(prs, report, df, category, prop, secondary, progress_cb, cof_y_range=cof_y_range)

        # Add branded footers to all content slides (skip title slide)
        for i, slide in enumerate(prs.slides):
            if i == 0:  # title slide has its own branding
                continue
            cls._add_footer(slide, logo_path, i + 1)

        # Save
        cat_slug = category.lower()
        filename = f"{project_name}_{cat_slug}_comparison.pptx"
        output_path = os.path.join(project_dir, filename)
        prs.save(output_path)
        return output_path

    # ------------------------------------------------------------------
    # Data helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _count_groups(report, category):
        prop = CATEGORY_PROP_MAP[category]
        seen = set()
        for f in report.files:
            val = getattr(f, prop, None)
            if val and str(val).strip():
                seen.add(str(val).strip())
        return max(len(seen), 1)

    @staticmethod
    def _build_summary_df(report, category, prop):
        """Build summary DataFrame with categorical columns."""
        rows = []
        for f in sorted(report.files, key=lambda x: x.file_basename.lower()):
            row = f.final_row.iloc[0].to_dict()
            row['Wafer'] = f.wafer_type or ''
            row['Pad'] = f.pad_type or ''
            row['Slurry'] = f.slurry_type or ''
            row['Conditioner'] = f.conditioner_disk_type or ''
            row['_raw_file'] = f  # keep reference for time traces
            rows.append(row)
        return pd.DataFrame(rows)

    @staticmethod
    def _get_groups(df, category):
        """Return dict of {group_name: sub_dataframe}, excluding unassigned."""
        valid = df[df[category].astype(str).str.strip() != '']
        groups = {}
        for name, gdf in valid.groupby(category):
            groups[str(name)] = gdf
        return groups

    @staticmethod
    def _get_psi_values(df):
        """Return sorted unique Pressure PSI values for series coloring."""
        if 'Pressure PSI' not in df.columns:
            return []
        vals = df['Pressure PSI'].dropna().unique()
        return sorted(vals)

    # ------------------------------------------------------------------
    # Slide helpers
    # ------------------------------------------------------------------

    @classmethod
    def _blank_slide(cls, prs):
        """Add a blank slide with branded background."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = SLIDE_BG
        return slide

    @classmethod
    def _add_title_bar(cls, slide, text):
        """Add a styled title bar with red accent line at the top of a slide."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0), TITLE_BAR_TOP,
            SLIDE_W, TITLE_BAR_H,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = TITLE_BAR_COLOR
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.name = 'Segoe UI'
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

        # Red accent line beneath title bar (echoes logo arc)
        accent = slide.shapes.add_shape(
            1,  # RECTANGLE
            Inches(0), Emu(int(TITLE_BAR_TOP + TITLE_BAR_H)),
            SLIDE_W, Inches(0.04),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT_COLOR
        accent.line.fill.background()

    @classmethod
    def _chart_positions(cls, n_charts, batch_start=0):
        """Return list of (left, top, width, height) for n_charts side-by-side."""
        n = min(n_charts, MAX_CHARTS_PER_SLIDE)
        usable_w = SLIDE_W - 2 * CHART_MARGIN
        chart_w = Emu(int(usable_w / n))
        positions = []
        for i in range(n):
            left = Emu(int(CHART_MARGIN + i * chart_w))
            positions.append((left, CHART_TOP, chart_w, CHART_H))
        return positions

    @classmethod
    def _format_chart_series(cls, chart, psi_values):
        """Apply colors and markers to chart series."""
        for i, series in enumerate(chart.series):
            color = SERIES_COLORS[i % len(SERIES_COLORS)]
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2)
            series.marker.style = XL_MARKER_STYLE.CIRCLE
            series.marker.size = 7
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color

    @classmethod
    def _format_bar_series(cls, chart, n_groups):
        """Apply colors to bar chart series."""
        for i, series in enumerate(chart.series):
            color = SERIES_COLORS[i % len(SERIES_COLORS)]
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color

    @classmethod
    def _set_log_scale(cls, axis, base=10):
        """Set logarithmic scale on a chart axis via XML."""
        scaling = axis._element.find(qn('c:scaling'))
        if scaling is not None:
            existing = scaling.find(qn('c:logBase'))
            if existing is not None:
                scaling.remove(existing)
            log_elem = etree.SubElement(scaling, qn('c:logBase'))
            log_elem.set('val', str(base))

    @staticmethod
    def _enable_gridlines(axis):
        """Enable major gridlines on a chart axis using the native API."""
        axis.has_major_gridlines = True
        axis.major_gridlines.format.line.color.rgb = RGBColor(0x3d, 0x3d, 0x3d)
        axis.major_gridlines.format.line.width = Pt(0.75)

    @classmethod
    def _add_group_label(cls, slide, text, left, top, width):
        """Add a text label above a chart identifying the group."""
        txbox = slide.shapes.add_textbox(left, Emu(int(top - Inches(0.25))), width, Inches(0.25))
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = 'Segoe UI'
        p.font.color.rgb = TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_footer(cls, slide, logo_path, slide_num=None):
        """Add branded footer with logo, company name, and slide number."""
        footer_y = Inches(7.1)

        # Separator line
        line = slide.shapes.add_shape(
            1, Inches(0.3), Inches(7.0),
            Emu(int(SLIDE_W - Inches(0.6))), Inches(0.01),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = FOOTER_LINE
        line.line.fill.background()

        # Logo (bottom-left)
        if logo_path and os.path.isfile(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.3), footer_y, height=Inches(0.35))

        # Company name (center)
        txbox = slide.shapes.add_textbox(Inches(4), footer_y, Inches(5), Inches(0.35))
        p = txbox.text_frame.paragraphs[0]
        p.text = "Araca Insights\u00ae"
        p.font.size = Pt(8)
        p.font.name = 'Segoe UI'
        p.font.color.rgb = FOOTER_TEXT
        p.alignment = PP_ALIGN.CENTER

        # Slide number (bottom-right)
        if slide_num is not None:
            txbox2 = slide.shapes.add_textbox(
                Emu(int(SLIDE_W - Inches(1))), footer_y,
                Inches(0.7), Inches(0.35),
            )
            p2 = txbox2.text_frame.paragraphs[0]
            p2.text = str(slide_num)
            p2.font.size = Pt(8)
            p2.font.name = 'Segoe UI'
            p2.font.color.rgb = FOOTER_TEXT
            p2.alignment = PP_ALIGN.RIGHT

    # ------------------------------------------------------------------
    # Slide sections
    # ------------------------------------------------------------------

    @classmethod
    def _add_title_slide(cls, prs, project_name, category, config, logo_path=None):
        slide = cls._blank_slide(prs)

        # Logo (centered above title)
        if logo_path and os.path.isfile(logo_path):
            from PIL import Image
            with Image.open(logo_path) as img:
                w, h = img.size
            logo_h = Inches(1.5)
            logo_w = Emu(int(logo_h * w / h))
            logo_left = Emu(int((SLIDE_W - logo_w) / 2))
            slide.shapes.add_picture(logo_path, logo_left, Inches(0.6), height=logo_h)

        # Main title
        title_top = Inches(2.3) if logo_path else Inches(2)
        txbox = slide.shapes.add_textbox(
            Inches(1), title_top, Inches(11), Inches(2)
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{category} vs {category} Comparison Report"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = 'Segoe UI Light'
        p.font.color.rgb = TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        p2 = tf.add_paragraph()
        p2.text = f"Project: {project_name}"
        p2.font.size = Pt(18)
        p2.font.name = 'Segoe UI'
        p2.font.color.rgb = TEXT_SECONDARY
        p2.alignment = PP_ALIGN.CENTER

        from datetime import date
        p3 = tf.add_paragraph()
        p3.text = date.today().strftime("%B %d, %Y")
        p3.font.size = Pt(14)
        p3.font.name = 'Segoe UI'
        p3.font.color.rgb = TEXT_MUTED
        p3.alignment = PP_ALIGN.CENTER

        # Red accent line below title text
        accent = slide.shapes.add_shape(
            1, Inches(4), Inches(4.4), Inches(5), Inches(0.03),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT_COLOR
        accent.line.fill.background()

        # Footer branding on title slide
        txbox_f = slide.shapes.add_textbox(Inches(4), Inches(6.8), Inches(5), Inches(0.4))
        pf = txbox_f.text_frame.paragraphs[0]
        pf.text = "Araca Insights\u00ae"
        pf.font.size = Pt(10)
        pf.font.name = 'Segoe UI'
        pf.font.color.rgb = FOOTER_TEXT
        pf.alignment = PP_ALIGN.CENTER

    @classmethod
    def _add_toc_slide(cls, prs, config, group_names):
        slide = cls._blank_slide(prs)
        cls._add_title_bar(slide, "Table of Contents")

        lines = [
            "Executive Summary",
        ]
        if config.get('selected_correlations', []):
            lines.append("Correlation Graphs")
        if config.get('include_summary_tables'):
            lines.append("Summary Tables")
        if config.get('include_time_traces'):
            lines.append("Appendix: Time Traces")

        txbox = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(11), Inches(5)
        )
        tf = txbox.text_frame
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{i + 1:02d}.  {line}"
            p.font.size = Pt(18)
            p.font.name = 'Segoe UI'
            p.font.color.rgb = TEXT_PRIMARY
            p.space_after = Pt(12)

    @classmethod
    def _add_executive_summary(cls, prs, df, groups, group_names, metrics, category, secondary, progress_cb):
        """Add clustered bar chart slides comparing groups across a secondary category."""
        if not metrics:
            return

        # Get unique non-empty values of the secondary category for x-axis
        sec_values = sorted(
            v for v in df[secondary].astype(str).str.strip().unique() if v
        )

        for batch_start in range(0, len(metrics), MAX_CHARTS_PER_SLIDE):
            batch = metrics[batch_start:batch_start + MAX_CHARTS_PER_SLIDE]
            slide = cls._blank_slide(prs)
            cls._add_title_bar(slide, "Executive Summary")

            positions = cls._chart_positions(len(batch))
            for i, metric in enumerate(batch):
                left, top, width, height = positions[i]
                label = SELECTION_LABELS.get(metric, metric)

                chart_data = CategoryChartData()

                if sec_values:
                    # Clustered: secondary values as x-axis, groups as series
                    chart_data.categories = sec_values
                    for gn in group_names:
                        gdf = groups[gn]
                        values = []
                        for sv in sec_values:
                            subset = gdf[gdf[secondary].astype(str).str.strip() == sv]
                            if subset.empty or metric not in subset.columns:
                                values.append(0)
                            else:
                                val = subset[metric].mean()
                                values.append(float(val) if pd.notna(val) else 0)
                        chart_data.add_series(gn, values)
                else:
                    # Fallback: group names as categories, single series
                    chart_data.categories = group_names
                    means = []
                    for gn in group_names:
                        gdf = groups[gn]
                        val = gdf[metric].mean() if metric in gdf.columns else 0
                        means.append(float(val) if pd.notna(val) else 0)
                    chart_data.add_series('Mean', means)

                chart_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED,
                    left, top, width, height, chart_data
                )
                chart = chart_frame.chart
                chart.has_title = True
                chart.chart_title.text_frame.paragraphs[0].text = label
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY

                # Style axes for dark theme
                chart.category_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
                chart.value_axis.tick_labels.font.color.rgb = TEXT_SECONDARY

                if sec_values:
                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False
                    chart.legend.font.color.rgb = TEXT_SECONDARY
                    cls._format_bar_series(chart, len(group_names))
                else:
                    chart.has_legend = False
                    series = chart.series[0]
                    for pt_idx in range(len(group_names)):
                        point = series.points[pt_idx]
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = SERIES_COLORS[pt_idx % len(SERIES_COLORS)]

            progress_cb("Creating executive summary...")

    @classmethod
    def _build_scatter_slide(cls, slide, group_names, groups, secondary, sv,
                             x_col, y_col, x_label, y_label, tx, ty,
                             psi_values, batch_start, batch_end,
                             y_min=None, y_max=None):
        """Build scatter charts on a slide for a batch of groups, filtered by secondary value.

        Returns True if any chart was added.
        """
        batch = group_names[batch_start:batch_end]
        positions = cls._chart_positions(len(batch))
        any_added = False

        for i, gn in enumerate(batch):
            left, top, width, height = positions[i]
            gdf = groups[gn]
            # Filter by secondary category value
            if sv:
                gdf = gdf[gdf[secondary].astype(str).str.strip() == sv]
            if gdf.empty:
                continue
            cls._add_group_label(slide, gn, left, top, width)

            # Compute transformed values
            work = gdf.copy()
            x_vals = cls._apply_transform(work, x_col, tx)
            y_vals = cls._apply_transform(work, y_col, ty)

            # Filter invalid
            mask = x_vals.notna() & y_vals.notna()
            if tx == 'log_axis':
                mask = mask & (x_vals > 0)
            if ty == 'log_axis':
                mask = mask & (y_vals > 0)
            mask = mask & np.isfinite(x_vals.fillna(0)) & np.isfinite(y_vals.fillna(0))

            x_clean = x_vals[mask]
            y_clean = y_vals[mask]
            work_clean = work[mask]

            chart_data = XyChartData()
            if psi_values and 'Pressure PSI' in work_clean.columns:
                for psi in psi_values:
                    sub_mask = work_clean['Pressure PSI'] == psi
                    xp = x_clean[sub_mask]
                    yp = y_clean[sub_mask]
                    if xp.empty:
                        continue
                    series = chart_data.add_series(f'{psi} PSI')
                    for xv, yv in zip(xp, yp):
                        series.add_data_point(float(xv), float(yv))
            else:
                if not x_clean.empty:
                    series = chart_data.add_series('All')
                    for xv, yv in zip(x_clean, y_clean):
                        series.add_data_point(float(xv), float(yv))

            if not any(len(s) > 0 for s in chart_data._series):
                continue

            chart_frame = slide.shapes.add_chart(
                XL_CHART_TYPE.XY_SCATTER_LINES,
                left, top, width, height, chart_data
            )
            chart = chart_frame.chart
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = TEXT_SECONDARY
            chart.has_title = False

            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_label
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = TEXT_SECONDARY
            chart.category_axis.has_title = True
            chart.category_axis.axis_title.text_frame.paragraphs[0].text = x_label
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
            chart.category_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = TEXT_SECONDARY

            chart.category_axis.tick_labels.font.size = Pt(7)
            chart.category_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
            chart.value_axis.tick_labels.font.size = Pt(7)
            chart.value_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
            chart.category_axis.tick_labels.number_format = '0.####'
            chart.value_axis.tick_labels.number_format = '0.####'

            cls._format_chart_series(chart, psi_values)

            if tx == 'log_axis':
                cls._set_log_scale(chart.category_axis)
            if ty == 'log_axis':
                cls._set_log_scale(chart.value_axis)
                chart.category_axis.crosses = XL_AXIS_CROSSES.MINIMUM
                chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

            if y_min is not None:
                chart.value_axis.minimum_scale = y_min
            if y_max is not None:
                chart.value_axis.maximum_scale = y_max

            cls._enable_gridlines(chart.value_axis)

            any_added = True

        return any_added

    @classmethod
    def _add_correlation_slides(cls, prs, df, groups, group_names, corr_configs, psi_values, category, secondary, progress_cb):
        """Add correlation graph slides, split by secondary category."""
        sec_values = sorted(
            v for v in df[secondary].astype(str).str.strip().unique() if v
        ) or [None]

        for entry in corr_configs:
            ci = entry['index'] if isinstance(entry, dict) else entry
            y_min = entry.get('y_min') if isinstance(entry, dict) else None
            y_max = entry.get('y_max') if isinstance(entry, dict) else None
            if ci >= len(PPTX_CORR_GRAPHS):
                continue
            graph_id, x_col, y_col, x_label, y_label, tx, ty = PPTX_CORR_GRAPHS[ci]

            for sv in sec_values:
                title = f"{y_label} \u2013 {sv}" if sv else f"{x_label} vs {y_label}"

                for batch_start in range(0, len(group_names), MAX_CHARTS_PER_SLIDE):
                    batch_end = batch_start + MAX_CHARTS_PER_SLIDE
                    slide = cls._blank_slide(prs)
                    cls._add_title_bar(slide, title)
                    cls._build_scatter_slide(
                        slide, group_names, groups, secondary, sv,
                        x_col, y_col, x_label, y_label, tx, ty,
                        psi_values, batch_start, batch_end,
                        y_min=y_min, y_max=y_max,
                    )

                progress_cb(f"Creating {graph_id} chart...")

    @classmethod
    def _apply_transform(cls, df, col, transform):
        """Apply data transform, returning a Series."""
        # Computed column: Directivity = Var Fy / Var Fz
        if col == '_directivity':
            if 'Var Fy' not in df.columns or 'Var Fz' not in df.columns:
                return pd.Series([np.nan] * len(df), index=df.index)
            var_fy = df['Var Fy'].copy()
            var_fz = df['Var Fz'].replace(0, np.nan)
            return var_fy / var_fz

        if col not in df.columns:
            return pd.Series([np.nan] * len(df), index=df.index)

        vals = df[col].copy()
        if transform == 'inv_kelvin':
            temp_k = vals + 273.15
            temp_k = temp_k.replace(0, np.nan)
            return 1.0 / temp_k
        elif transform == 'ln':
            vals = vals.replace(0, np.nan)
            return np.log(vals)
        else:
            return vals

    @classmethod
    def _add_summary_tables(cls, prs, df, groups, group_names, features, category, secondary, progress_cb):
        """Add one summary table slide per (group × secondary value) combination."""
        if not features:
            return

        sec_values = sorted(
            v for v in df[secondary].astype(str).str.strip().unique() if v
        ) or [None]

        col_labels = [SELECTION_LABELS.get(f, f) for f in features]

        for gn in group_names:
            gdf = groups[gn]
            for sv in sec_values:
                # Filter by secondary category
                if sv:
                    sub = gdf[gdf[secondary].astype(str).str.strip() == sv]
                else:
                    sub = gdf
                if sub.empty:
                    continue

                # Sort by Pressure PSI then Mean Velocity
                sort_cols = [c for c in ['Pressure PSI', 'Mean Velocity'] if c in sub.columns]
                if sort_cols:
                    sub = sub.sort_values(sort_cols)

                # Compute Directivity if selected
                has_directivity = 'Directivity' in features
                if has_directivity and 'Var Fy' in sub.columns and 'Var Fz' in sub.columns:
                    var_fz = sub['Var Fz'].replace(0, np.nan)
                    directivity_vals = sub['Var Fy'] / var_fz
                else:
                    directivity_vals = pd.Series([np.nan] * len(sub), index=sub.index)

                title = f"Summary Table \u2013 {gn} \u2013 {sv}" if sv else f"Summary Table \u2013 {gn}"
                slide = cls._blank_slide(prs)
                cls._add_title_bar(slide, title)

                n_rows = len(sub) + 1  # +1 for header
                n_cols = len(features)

                table_left = Inches(0.3)
                table_top = Inches(1.3)
                table_w = Emu(int(SLIDE_W - Inches(0.6)))
                table_h = Inches(min(5.5, 0.35 * n_rows + 0.35))

                table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h)
                table = table_shape.table

                # Header row
                for c, label in enumerate(col_labels):
                    cell = table.cell(0, c)
                    cell.text = label
                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for run in p.runs:
                            run.font.size = Pt(9)
                            run.font.bold = True
                            run.font.name = 'Segoe UI'
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TITLE_BAR_COLOR

                # Data rows (alternating colors for readability)
                for r, (idx, row) in enumerate(sub.iterrows(), start=1):
                    row_color = TABLE_ROW_ALT if r % 2 == 0 else TABLE_ROW_WHITE
                    for c, feat in enumerate(features):
                        if feat == 'Directivity':
                            val = directivity_vals.get(idx, np.nan)
                        else:
                            val = row.get(feat, '')
                        if isinstance(val, float) and pd.notna(val):
                            if val != 0 and abs(val) < 0.001:
                                text = f"{val:.3e}"
                            elif abs(val) < 10:
                                text = f"{val:.3f}"
                            elif abs(val) < 1000:
                                text = f"{val:.1f}"
                            else:
                                text = f"{val:,.0f}"
                        else:
                            text = str(val) if pd.notna(val) else ''
                        cell = table.cell(r, c)
                        cell.text = text
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = row_color
                        for p in cell.text_frame.paragraphs:
                            p.alignment = PP_ALIGN.CENTER
                            for run in p.runs:
                                run.font.size = Pt(9)
                                run.font.name = 'Segoe UI'
                                run.font.color.rgb = TEXT_PRIMARY

                progress_cb(f"Creating summary table for {gn} \u2013 {sv}..." if sv else f"Creating summary table for {gn}...")

    @classmethod
    def _add_time_traces(cls, prs, report, df, category, prop, secondary, progress_cb, *, cof_y_range=None):
        """Add appendix with one slide per file showing time-series charts."""
        sorted_files = sorted(report.files, key=lambda f: f.file_basename.lower())
        secondary_prop = CATEGORY_PROP_MAP.get(secondary, '')

        for raw_file in sorted_files:
            slide = cls._blank_slide(prs)

            # Build descriptive title from file properties
            primary_val = getattr(raw_file, prop, '') or ''
            secondary_val = getattr(raw_file, secondary_prop, '') or '' if secondary_prop else ''
            psi = raw_file.pressure_psi
            try:
                velocity = raw_file.final_row['Mean Velocity'].iloc[0]
            except (KeyError, IndexError):
                velocity = 0

            parts = ['Time Traces']
            if primary_val:
                parts.append(primary_val)
            if psi:
                parts.append(f"{round(psi, 1)} PSI")
            if velocity:
                parts.append(f"{round(velocity, 1)} m/s")
            if secondary_val:
                parts.append(secondary_val)
            title = ' \u2013 '.join(parts)
            cls._add_title_bar(slide, title)

            tpf = raw_file.total_per_frame
            if tpf is None or tpf.empty:
                progress_cb(f"Creating time trace for {raw_file.file_basename}...")
                continue

            time_col = 'time (s)' if 'time (s)' in tpf.columns else None
            if not time_col:
                progress_cb(f"Creating time trace for {raw_file.file_basename}...")
                continue

            positions = cls._chart_positions(len(TIME_TRACE_METRICS))
            for i, (chart_title, series_defs) in enumerate(TIME_TRACE_METRICS):
                # Skip chart if none of its columns exist
                available = [(cn, sl) for cn, sl in series_defs if cn in tpf.columns]
                if not available:
                    continue

                left, top, width, height = positions[i]
                chart_data = XyChartData()
                has_data = False

                for col_name, series_label in available:
                    series = chart_data.add_series(series_label)
                    data = tpf[[time_col, col_name]].dropna()
                    if len(data) > 500:
                        step = max(1, len(data) // 500)
                        data = data.iloc[::step]
                    for _, row in data.iterrows():
                        series.add_data_point(float(row[time_col]), float(row[col_name]))
                    if len(series) > 0:
                        has_data = True

                if not has_data:
                    continue

                chart_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
                    left, top, width, height, chart_data
                )
                chart = chart_frame.chart
                chart.has_legend = len(available) > 1
                if chart.has_legend:
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False
                    chart.legend.font.size = Pt(8)
                    chart.legend.font.color.rgb = TEXT_SECONDARY
                chart.has_title = True
                chart.chart_title.text_frame.paragraphs[0].text = chart_title
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = TEXT_PRIMARY

                # Axis – force x-axis labels to bottom even when data goes negative
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.paragraphs[0].text = 'Time (s)'
                chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(9)
                chart.category_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = TEXT_SECONDARY
                chart.category_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
                chart.value_axis.tick_labels.font.color.rgb = TEXT_SECONDARY
                chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

                # Apply user-specified COF Y-axis range
                if chart_title == 'COF' and cof_y_range:
                    y_min, y_max = cof_y_range
                    if y_min is not None:
                        chart.value_axis.minimum_scale = y_min
                    if y_max is not None:
                        chart.value_axis.maximum_scale = y_max

                for si, (_, _) in enumerate(available):
                    s = chart.series[si]
                    s.format.line.color.rgb = SERIES_COLORS[si % len(SERIES_COLORS)]
                    s.format.line.width = Pt(1.5)

            progress_cb(f"Creating time trace for {raw_file.file_basename}...")
